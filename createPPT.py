from pptx import Presentation

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Use layout for title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.placeholders[0]
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

def main():
    prs = Presentation()

    create_slide(prs, "Unclear Roles and Blame Shifting",
                 "Revisit team roles in Scrum framework.\nEnsure understanding of responsibilities.\nFoster culture of accountability and collaboration.")
    
    create_slide(prs, "Lack of Code Reviews and Test Coverage Reports",
                 "Implement mandatory code reviews for all changes.\nUse GitHub pull requests for code reviews.\nRequire reviews before deployment.")
    
    create_slide(prs, "Delays Due to Dependencies and Unmaintained Backlog",
                 "Improve backlog grooming practices.\nBreak down user stories into smaller tasks.\nRegularly review and prioritize the backlog.")
    
    create_slide(prs, "Integration of Code Reviews into CI/CD Pipeline",
                 "Configure Jenkins to trigger code analysis.\nUse tools like SonarQube or CodeClimate for automated code reviews.\nRequire passing code quality checks before deployment.")
    
    create_slide(prs, "Dependency Management and Build Isolation",
                 "Use NuGet for package management.\nDefine clear dependencies and versioning in project files.\nImplement build isolation techniques.")

    prs.save("Agile_Development_Solutions.pptx")

if __name__ == "__main__":
    main()

